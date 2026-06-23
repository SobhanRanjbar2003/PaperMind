"""
Builds the final .pptx file from the normalized slide plan produced by slide_planner.

Design principles:
  - Full RTL support (Persian/Arabic text): alignment, XML rtl flag, complex-script font.
  - One dominant color per palette (~60% weight) with 1–2 supporting tones and one accent.
  - A repeating visual motif (numbered circles) carried across all content slides.
  - Six layout types: image-right, image-left, bullets-only, quote, two-column, stat.
  - Decorative shapes replace missing images — never empty white space.
  - palette["style"] drives the decorative motif variant (circles, wave, leaf, bold, minimal).
"""

import io
import logging

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from app.config import settings

logger = logging.getLogger(__name__)

# ── Layout constants ──────────────────────────────────────────────────────────
MARGIN = Inches(0.6)
GAP = Inches(0.45)
IMG_COL_W = Inches(5.1)
FONT_FA = settings.presentation_font_fa


# ── Color helpers ─────────────────────────────────────────────────────────────

def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str.lstrip("#"))


def _luminance(hex_str: str) -> float:
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return 0.299 * r + 0.587 * g + 0.114 * b


def _contrast_text(bg_hex: str) -> str:
    """Return white or dark text for readable contrast over bg_hex."""
    return "FFFFFF" if _luminance(bg_hex) < 150 else "1A1A1A"


def _best_text_color(bg_hex: str, preferred_hex: str, min_diff: float = 90.0) -> str:
    if abs(_luminance(bg_hex) - _luminance(preferred_hex)) >= min_diff:
        return preferred_hex
    return "F2F2F2" if _luminance(bg_hex) < 150 else "1A1A1A"


def _pick_contrasting(bg_hex: str, candidates: list[str], min_diff: float = 70.0) -> str:
    for c in candidates:
        if abs(_luminance(bg_hex) - _luminance(c)) >= min_diff:
            return c
    return "F2F2F2" if _luminance(bg_hex) < 150 else "1A1A1A"


# ── Shape / text primitives ───────────────────────────────────────────────────

def _set_rtl(paragraph) -> None:
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def _style_run(run, size_pt: float, color_hex: str, bold: bool = False, italic: bool = False) -> None:
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _rgb(color_hex)
    run.font.name = FONT_FA
    rPr = run._r.get_or_add_rPr()
    for tag, successors in (
        ("a:ea", ("a:cs", "a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")),
        ("a:cs", ("a:sym", "a:hlinkClick", "a:hlinkMouseOver", "a:rtl", "a:extLst")),
    ):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.insert_element_before(el, *successors)
        el.set("typeface", FONT_FA)


def _set_alpha(shape, alpha_pct: int) -> None:
    solidFill = shape.fill.fore_color._xFill
    srgbClr = solidFill.find(qn("a:srgbClr"))
    if srgbClr is not None:
        alpha = srgbClr.makeelement(qn("a:alpha"), {})
        alpha.set("val", str(int(alpha_pct * 1000)))
        srgbClr.append(alpha)


def _no_line(shape) -> None:
    shape.line.fill.background()
    try:
        shape.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass


def _solid(shape, hex_color: str, alpha_pct: int | None = None) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(hex_color)
    if alpha_pct is not None and alpha_pct < 100:
        _set_alpha(shape, alpha_pct)


def _set_bg(slide, hex_color: str) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(hex_color)


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    return tb, tf


# ── Image insertion ───────────────────────────────────────────────────────────

def _add_cover_image(slide, image_bytes: bytes, left, top, width, height):
    stream = io.BytesIO(image_bytes)
    try:
        with Image.open(io.BytesIO(image_bytes)) as img:
            iw, ih = img.size
    except Exception:  # noqa: BLE001
        iw, ih = 1, 1

    pic = slide.shapes.add_picture(stream, left, top, width=width, height=height)
    if iw > 0 and ih > 0:
        img_ratio = iw / ih
        box_ratio = width / height
        if img_ratio > box_ratio:
            crop = max(0.0, min(0.49, (1 - box_ratio / img_ratio) / 2))
            pic.crop_left = crop
            pic.crop_right = crop
        elif img_ratio < box_ratio:
            crop = max(0.0, min(0.49, (1 - img_ratio / box_ratio) / 2))
            pic.crop_top = crop
            pic.crop_bottom = crop
    return pic


# ── Decorative motifs ─────────────────────────────────────────────────────────

def _motif_circles(slide, x_anchor, y_anchor, color_a: str, color_b: str, scale: float = 1.0) -> None:
    """Overlapping translucent circles — default motif."""
    d1 = Inches(3.2 * scale)
    d2 = Inches(2.0 * scale)
    c1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, x_anchor - d1 // 2, y_anchor - d1 // 2, d1, d1)
    _solid(c1, color_a, alpha_pct=22)
    _no_line(c1)
    c2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        x_anchor - d2 // 2 + Inches(0.6 * scale),
        y_anchor - d2 // 2 + Inches(0.4 * scale),
        d2, d2,
    )
    _solid(c2, color_b, alpha_pct=35)
    _no_line(c2)


def _motif_wave(slide, x_anchor, y_anchor, color_a: str, color_b: str, scale: float = 1.0) -> None:
    """Stacked rounded rectangles giving a layered/wave feel."""
    w = Inches(3.5 * scale)
    h = Inches(1.0 * scale)
    for i, (col, alpha) in enumerate([(color_a, 18), (color_b, 28), (color_a, 15)]):
        r = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_anchor - w // 2,
            y_anchor - Inches(0.8 * scale) + Inches(i * 0.7 * scale),
            w, h,
        )
        _solid(r, col, alpha_pct=alpha)
        _no_line(r)


def _motif_leaf(slide, x_anchor, y_anchor, color_a: str, color_b: str, scale: float = 1.0) -> None:
    """Diamond shapes suggesting leaves or growth."""
    d = Inches(2.2 * scale)
    for offset_x, offset_y, col, alpha in [
        (0, 0, color_a, 20),
        (Inches(0.9 * scale), Inches(0.5 * scale), color_b, 30),
    ]:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            x_anchor - d // 2 + offset_x,
            y_anchor - d // 2 + offset_y,
            d, d,
        )
        _solid(shape, col, alpha_pct=alpha)
        _no_line(shape)


def _motif_bold(slide, x_anchor, y_anchor, color_a: str, color_b: str, scale: float = 1.0) -> None:
    """Bold solid square + circle — high-energy style."""
    d = Inches(2.4 * scale)
    sq = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x_anchor - d // 2, y_anchor - d // 2, d, d
    )
    _solid(sq, color_b, alpha_pct=30)
    _no_line(sq)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x_anchor - Inches(0.8 * scale), y_anchor - Inches(0.8 * scale),
        Inches(1.6 * scale), Inches(1.6 * scale),
    )
    _solid(circle, color_a, alpha_pct=45)
    _no_line(circle)


def _motif_minimal(slide, x_anchor, y_anchor, color_a: str, color_b: str, scale: float = 1.0) -> None:
    """Single thin-bordered circle — clean and minimal."""
    d = Inches(2.8 * scale)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x_anchor - d // 2, y_anchor - d // 2, d, d
    )
    _solid(circle, color_a, alpha_pct=12)
    circle.line.color.rgb = _rgb(color_b)
    circle.line.width = Pt(1.5)


_MOTIF_FN = {
    "circles": _motif_circles,
    "wave": _motif_wave,
    "leaf": _motif_leaf,
    "bold": _motif_bold,
    "minimal": _motif_minimal,
    "rounded": _motif_circles,  # fallback alias
}


def _add_motif(slide, style: str, x_anchor, y_anchor, color_a: str, color_b: str, scale: float = 1.0) -> None:
    fn = _MOTIF_FN.get(style, _motif_circles)
    fn(slide, x_anchor, y_anchor, color_a, color_b, scale)


# ── Bullet rows (numbered circles) ───────────────────────────────────────────

def _add_bullet_rows(
    slide,
    bullets: list[str],
    x, y, width, height,
    circle_color: str,
    text_color: str,
    font_size: float = 16,
) -> None:
    n = len(bullets)
    if n == 0:
        return

    ideal_row_h = Inches(1.3)
    row_h = Emu(min(int(ideal_row_h), int(height / n)))
    total_h = Emu(int(row_h * n))
    y = Emu(int(y + max(0, (height - total_h) / 2)))

    circle_d = Inches(0.5)
    gap = Inches(0.18)

    for i, bullet in enumerate(bullets):
        row_top = Emu(int(y + i * row_h))

        # Numbered circle
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Emu(int(x + width - circle_d)),
            Emu(int(row_top + (row_h - circle_d) / 2)),
            circle_d, circle_d,
        )
        _solid(circle, circle_color)
        _no_line(circle)
        ctf = circle.text_frame
        ctf.word_wrap = False
        ctf.margin_left = ctf.margin_right = ctf.margin_top = ctf.margin_bottom = 0
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = str(i + 1)
        _style_run(cr, 15, _contrast_text(circle_color), bold=True)

        # Bullet text
        text_w = Emu(int(width - circle_d - gap))
        _tb, tf = _textbox(slide, x, row_top, text_w, row_h)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        _set_rtl(p)
        run = p.add_run()
        run.text = bullet
        _style_run(run, font_size, text_color)


def _add_speaker_notes(slide, slide_data: dict) -> None:
    notes = slide_data.get("speaker_notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


# ── Slide title (shared) ──────────────────────────────────────────────────────

def _add_content_title(slide, x, width, text: str, color: str) -> Emu:
    top = Inches(0.55)
    height = Inches(1.15)
    _tb, tf = _textbox(slide, x, top, width, height)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_rtl(p)
    run = p.add_run()
    run.text = text
    _style_run(run, 32, color, bold=True)
    return Emu(int(top + height))


# ── Title slide ───────────────────────────────────────────────────────────────

def _add_title_slide(prs, layout, plan: dict, cover_image: bytes | None,
                     primary: str, secondary: str, accent: str, style: str) -> None:
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height

    if cover_image:
        _add_cover_image(slide, cover_image, 0, 0, sw, sh)
        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, sw, sh)
        _solid(overlay, primary, alpha_pct=72)
        _no_line(overlay)
    else:
        _set_bg(slide, primary)
        _add_motif(slide, style, sw - Inches(2.2), sh - Inches(1.8), secondary, accent, scale=1.3)
        _add_motif(slide, style, Inches(1.4), Inches(1.2), accent, secondary, scale=0.8)

    text_color = _contrast_text(primary)

    _tb, tf = _textbox(slide, Inches(1.0), sh / 2 - Inches(1.1), sw - Inches(2.0), Inches(1.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_rtl(p)
    run = p.add_run()
    run.text = plan["presentation_title"]
    _style_run(run, 44, text_color, bold=True)

    if plan.get("subtitle"):
        _tb2, tf2 = _textbox(slide, Inches(1.5), sh / 2 + Inches(0.55), sw - Inches(3.0), Inches(0.9))
        p2 = tf2.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        _set_rtl(p2)
        run2 = p2.add_run()
        run2.text = plan["subtitle"]
        _style_run(run2, 20, _best_text_color(primary, secondary))


# ── Content slides ────────────────────────────────────────────────────────────

def _add_image_variant_slide(prs, layout, slide_data: dict,
                              primary: str, secondary: str, accent: str,
                              style: str, image_bytes: bytes | None, image_left: bool) -> None:
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    bg = slide_data.get("_bg", "FFFFFF")
    _set_bg(slide, bg)

    if image_left:
        img_x, text_x = 0, Emu(int(IMG_COL_W + GAP))
    else:
        img_x, text_x = Emu(int(sw - IMG_COL_W)), MARGIN

    text_w = Emu(int(sw - IMG_COL_W - GAP - MARGIN))

    if image_bytes:
        _add_cover_image(slide, image_bytes, img_x, 0, IMG_COL_W, sh)
    else:
        block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_x, 0, IMG_COL_W, sh)
        _solid(block, secondary)
        _no_line(block)
        _add_motif(slide, style, Emu(int(img_x + IMG_COL_W / 2)), Emu(int(sh / 2)), primary, accent, scale=1.1)

    title_bottom = _add_content_title(slide, text_x, text_w, slide_data["title"], "1A1A1A")
    bullets_top = Emu(int(title_bottom + Inches(0.25)))
    bullets_height = Emu(int(sh - bullets_top - Inches(0.55)))
    _add_bullet_rows(slide, slide_data["bullets"], text_x, bullets_top, text_w, bullets_height,
                     circle_color=primary, text_color="2B2B2B", font_size=16)
    _add_speaker_notes(slide, slide_data)


def _add_bullets_only_slide(prs, layout, slide_data: dict,
                             primary: str, secondary: str, accent: str, style: str) -> None:
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    bg = slide_data.get("_bg", "FFFFFF")
    _set_bg(slide, bg)

    # Subtle motif in top-right corner
    _add_motif(slide, style, sw - Inches(1.0), Inches(0.8), secondary, accent, scale=0.55)

    content_w = Emu(int(sw - 2 * MARGIN))
    title_bottom = _add_content_title(slide, MARGIN, content_w, slide_data["title"], "1A1A1A")
    bullets_top = Emu(int(title_bottom + Inches(0.35)))
    bullets_height = Emu(int(sh - bullets_top - Inches(0.6)))
    bullets_w = Emu(int(content_w * 0.62))
    bullets_x = Emu(int(MARGIN + (content_w - bullets_w) / 2))
    _add_bullet_rows(slide, slide_data["bullets"], bullets_x, bullets_top, bullets_w, bullets_height,
                     circle_color=primary, text_color="2B2B2B", font_size=18)
    _add_speaker_notes(slide, slide_data)


def _add_quote_slide(prs, layout, slide_data: dict,
                     primary: str, secondary: str, accent: str, style: str) -> None:
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    _set_bg(slide, primary)

    _add_motif(slide, style, Inches(1.3), sh - Inches(1.2), accent, secondary, scale=1.0)
    _add_motif(slide, style, sw - Inches(1.3), Inches(1.2), secondary, accent, scale=0.9)

    text_color = _contrast_text(primary)
    quote_text = (slide_data.get("bullets") or [slide_data["title"]])[0]

    # Opening quote mark
    _tb0, tf0 = _textbox(slide, Inches(1.0), sh / 2 - Inches(1.7), Inches(1.0), Inches(0.8))
    p0 = tf0.paragraphs[0]
    p0.alignment = PP_ALIGN.RIGHT
    r0 = p0.add_run()
    r0.text = "«"
    _style_run(r0, 72, _best_text_color(primary, accent), bold=True)

    _tb, tf = _textbox(slide, Inches(1.6), sh / 2 - Inches(1.2), sw - Inches(3.2), Inches(2.0))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_rtl(p)
    run = p.add_run()
    run.text = quote_text
    _style_run(run, 28, text_color, italic=True)

    _tb2, tf2 = _textbox(slide, Inches(2.0), sh / 2 + Inches(1.0), sw - Inches(4.0), Inches(0.8))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _set_rtl(p2)
    run2 = p2.add_run()
    run2.text = slide_data["title"]
    _style_run(run2, 18, _best_text_color(primary, accent), bold=True)

    _add_speaker_notes(slide, slide_data)


def _add_two_column_slide(prs, layout, slide_data: dict,
                           primary: str, secondary: str, accent: str, style: str) -> None:
    """
    Two-column comparison layout.
    Bullets split on " | " into left and right columns.
    If no separator found, first half goes left, second half right.
    """
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    bg = slide_data.get("_bg", "FFFFFF")
    _set_bg(slide, bg)

    content_w = Emu(int(sw - 2 * MARGIN))
    title_bottom = _add_content_title(slide, MARGIN, content_w, slide_data["title"], "1A1A1A")

    bullets = slide_data.get("bullets", [])
    left_bullets, right_bullets = [], []
    for b in bullets:
        if " | " in b:
            parts = b.split(" | ", 1)
            left_bullets.append(parts[0].strip())
            right_bullets.append(parts[1].strip())
        else:
            left_bullets.append(b)

    if not right_bullets:
        mid = max(1, len(left_bullets) // 2)
        right_bullets = left_bullets[mid:]
        left_bullets = left_bullets[:mid]

    col_gap = Inches(0.3)
    col_w = Emu(int((content_w - col_gap) / 2))
    bullets_top = Emu(int(title_bottom + Inches(0.25)))
    bullets_height = Emu(int(sh - bullets_top - Inches(0.6)))

    # Column headers
    for x_pos, label, col in [
        (MARGIN, "ستون راست", primary),
        (Emu(int(MARGIN + col_w + col_gap)), "ستون چپ", secondary),
    ]:
        header_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x_pos, bullets_top,
                                           col_w, Inches(0.42))
        _solid(header_bg, col)
        _no_line(header_bg)

    # Right column bullets
    _add_bullet_rows(slide, right_bullets, MARGIN,
                     Emu(int(bullets_top + Inches(0.5))),
                     col_w, Emu(int(bullets_height - Inches(0.5))),
                     circle_color=primary, text_color="2B2B2B", font_size=15)

    # Left column bullets
    _add_bullet_rows(slide, left_bullets,
                     Emu(int(MARGIN + col_w + col_gap)),
                     Emu(int(bullets_top + Inches(0.5))),
                     col_w, Emu(int(bullets_height - Inches(0.5))),
                     circle_color=_pick_contrasting(bg, [secondary, accent, primary]),
                     text_color="2B2B2B", font_size=15)

    _add_speaker_notes(slide, slide_data)


def _add_stat_slide(prs, layout, slide_data: dict,
                    primary: str, secondary: str, accent: str, style: str) -> None:
    """
    Stat/number callout layout — large centered numbers with labels below.
    Each bullet treated as "number — description" or just shown large.
    """
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    bg = slide_data.get("_bg", "FFFFFF")
    _set_bg(slide, bg)

    content_w = Emu(int(sw - 2 * MARGIN))
    title_bottom = _add_content_title(slide, MARGIN, content_w, slide_data["title"], "1A1A1A")

    bullets = slide_data.get("bullets", [])[:4]
    n = len(bullets) or 1
    card_w = Emu(int((content_w - Inches(0.25) * (n - 1)) / n))
    card_top = Emu(int(title_bottom + Inches(0.4)))
    card_h = Emu(int(sh - card_top - Inches(0.6)))

    colors_cycle = [primary, secondary, accent, primary]

    for i, bullet in enumerate(bullets):
        cx = Emu(int(MARGIN + i * (card_w + Inches(0.25))))

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, card_top, card_w, card_h)
        card_color = colors_cycle[i % len(colors_cycle)]
        _solid(card, card_color, alpha_pct=15)
        _no_line(card)

        # Split "number — label" if possible
        if "—" in bullet:
            number_part, label_part = bullet.split("—", 1)
        elif "-" in bullet and bullet.index("-") < 10:
            number_part, label_part = bullet.split("-", 1)
        else:
            number_part = bullet[:6]
            label_part = bullet[6:]

        # Large number
        _tb, tf = _textbox(slide, cx, card_top, card_w, Emu(int(card_h * 0.55)))
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = number_part.strip()
        _style_run(run, 42, _pick_contrasting(bg, [primary, secondary, "1A1A1A"]), bold=True)

        # Label below
        if label_part.strip():
            _tb2, tf2 = _textbox(slide, cx, Emu(int(card_top + card_h * 0.55)),
                                 card_w, Emu(int(card_h * 0.4)))
            tf2.vertical_anchor = MSO_ANCHOR.TOP
            p2 = tf2.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            _set_rtl(p2)
            run2 = p2.add_run()
            run2.text = label_part.strip()
            _style_run(run2, 14, "2B2B2B")

    _add_speaker_notes(slide, slide_data)


# ── Closing slide ─────────────────────────────────────────────────────────────

def _add_closing_slide(prs, layout, plan: dict,
                       primary: str, secondary: str, accent: str, style: str) -> None:
    slide = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    _set_bg(slide, primary)

    _add_motif(slide, style, sw - Inches(2.0), Inches(1.5), secondary, accent, scale=1.2)

    text_color = _contrast_text(primary)

    _tb, tf = _textbox(slide, Inches(1.0), Inches(0.9), sw - Inches(2.0), Inches(1.0))
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_rtl(p)
    run = p.add_run()
    run.text = "جمع‌بندی"
    _style_run(run, 38, text_color, bold=True)

    recap_titles = [s["title"] for s in plan["slides"][:6]]
    bullets_w = Emu(int((sw - Inches(2.0)) * 0.8))
    bullets_x = Emu(int(Inches(1.0) + ((sw - Inches(2.0)) - bullets_w) / 2))
    bullets_top = Inches(2.1)
    bullets_h = Emu(int(sh - bullets_top - Inches(1.1)))
    _add_bullet_rows(slide, recap_titles, bullets_x, bullets_top, bullets_w, bullets_h,
                     circle_color=_pick_contrasting(primary, [secondary, accent]),
                     text_color=text_color, font_size=18)

    _tb2, tf2 = _textbox(slide, Inches(1.5), sh - Inches(0.9), sw - Inches(3.0), Inches(0.6))
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    _set_rtl(p2)
    run2 = p2.add_run()
    run2.text = f"پایان ارائه — بر اساس خلاصه‌ی «{plan['presentation_title']}»"
    _style_run(run2, 14, _best_text_color(primary, accent), italic=True)


# ── Main entry point ──────────────────────────────────────────────────────────

def build_presentation(
    plan: dict,
    slide_images: dict[int, bytes],
    cover_image: bytes | None,
    output_path: str,
) -> str:
    """
    Build the final .pptx file from a normalized slide plan.

    Args:
        plan:         Output of slide_planner.build_slide_plan / _normalize_plan.
        slide_images: {slide_index: image_bytes} — empty dict if no images.
        cover_image:  Bytes for title slide background image, or None.
        output_path:  Absolute path to write the .pptx file.
    """
    prs = Presentation()
    prs.slide_width = Inches(settings.presentation_slide_width_in)
    prs.slide_height = Inches(settings.presentation_slide_height_in)
    layout = prs.slide_layouts[6]  # completely blank layout

    colors = plan["palette_colors"]
    primary = colors["primary"]
    secondary = colors["secondary"]
    accent = colors["accent"]
    bg = colors.get("bg", "FFFFFF")
    style = colors.get("style", "circles")

    prs.core_properties.title = plan["presentation_title"]
    prs.core_properties.author = "Book Summarizer"

    _add_title_slide(prs, layout, plan, cover_image, primary, secondary, accent, style)

    for idx, slide_data in enumerate(plan["slides"]):
        # Inject bg into slide_data for content builders
        slide_data["_bg"] = bg
        image_bytes = slide_images.get(idx)
        layout_kind = slide_data["layout"]

        if layout_kind == "image-right":
            _add_image_variant_slide(prs, layout, slide_data, primary, secondary, accent,
                                     style, image_bytes, image_left=False)
        elif layout_kind == "image-left":
            _add_image_variant_slide(prs, layout, slide_data, primary, secondary, accent,
                                     style, image_bytes, image_left=True)
        elif layout_kind == "quote":
            _add_quote_slide(prs, layout, slide_data, primary, secondary, accent, style)
        elif layout_kind == "two-column":
            _add_two_column_slide(prs, layout, slide_data, primary, secondary, accent, style)
        elif layout_kind == "stat":
            _add_stat_slide(prs, layout, slide_data, primary, secondary, accent, style)
        else:  # bullets-only or unknown
            _add_bullets_only_slide(prs, layout, slide_data, primary, secondary, accent, style)

    _add_closing_slide(prs, layout, plan, primary, secondary, accent, style)

    prs.save(output_path)
    logger.info("PPTX saved: %s", output_path)
    return output_path
